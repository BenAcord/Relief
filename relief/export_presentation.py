from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
import time

# Initialize the presentation
prs = Presentation()
killChainNames=["Recon", "Weapon.", "Delivery", "Exploit", "Install", "C2", "Actions"]
timestr = time.strftime("%Y%m%d-%H%M")
presentationName = "export/relief-CAMPAIGNNAME-" + timestr + ".pptx"

#--------
#      Default is 4:3 need to switch this to 16:9 or pick an option in settings.
#--------


# Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

# Title slide content
title.text = "Campaign Title"
subtitle.text = "Campaign date"

# Executive Summary Slide
bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Executive Summary'

tf = body_shape.text_frame
tf.text = 'Campaign Name'

p = tf.add_paragraph()
p.text = 'Summary of Business Impact'
p.level = 1

p = tf.add_paragraph()
p.text = '$ of lost something'
p.level = 2

p = tf.add_paragraph()
p.text = "% of time lost"
p.level = 2

p = tf.add_paragraph()
p.text = '# of widgets damaged'
p.level = 2

# Kill chain hotlinks
left = Inches(0.2)  # 0.93" centers this overall set of shapes
top = Inches(6.0)
width = Inches(1.25)
height = Inches(.75)

shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.text = 'Recon'

left = left + width - Inches(0.0)
width = Inches(1.20)  # chevrons need more width for visual balance, set to 2.0

for n in killChainNames:
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.text = '%s' % n
    left = left + width - Inches(0.0)



# Detail slide
# reset values
left = Inches(0)
top = Inches(0)
width = Inches(0)
height = Inches(0)

title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

shapes.title.text = 'Kill Chain <<NAME>> Details'

left = Inches(1.0)  # 0.93" centers this overall set of shapes
top = Inches(2.0)
width = Inches(1.75)
height = Inches(1.0)

img_path = 'static/images/diamond.png'

txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

tf.text = "This is text inside a textbox"

p = tf.add_paragraph()
p.text = "This is a second paragraph that's bold"
p.font.bold = True

p = tf.add_paragraph()
p.text = "This is a third paragraph that's big"
p.font.size = Pt(40)

# orient the diamond image
left = top = Inches(2.5)
pic = slide.shapes.add_picture(img_path, left, top)


# Write the slides to a new presentation file
print("[+]  Creating slide deck for campain <<CAMPAIGN NAME>> in ", presentationName)
prs.save(presentationName)