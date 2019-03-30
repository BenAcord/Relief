from docx import Document
from docx.shared import Inches
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.text import WD_TAB_ALIGNMENT
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches, Pt
from flask import abort, flash, render_template, jsonify, request, url_for, redirect, Blueprint
from relief import app, db
from relief.forms import UpdateSettingsForm, NewCampaign, AddKillChainPhase, DisplayCampaign, ListCampaigns, Home, Stats
from relief.models import Campaigns, ReliefConfig, CyberKillChains
import relief.settings
import json
from time import strftime, time
import datetime
import os


@app.route('/campaigns')
def campaignspage():
    return render_template('campaigns.html')

@app.route('/new/', methods=["GET", "POST"])
def newpage():
    #return render_template("new.html")
    return redirect(url_for('newcampaign'))

@app.route('/')
@app.route('/home', methods=["GET"])
def homepage():
    form = Home()
    return render_template("home.html", title='Home', form=form)

@app.route('/stats', methods=["GET"])
def statspage():
    form = Stats()
    return render_template('stats.html', title='Stats', form=form)

@app.route('/newcampaign/', methods=["GET", "POST"])
def newcampaign():
    form = NewCampaign()
    if form.validate_on_submit():
        new_campaign = Campaigns(campaign_name = form.campaignName.data, 
            report_date = form.reportDate.data, 
            summary_of_business_impact = form.summaryBusinessImpact.data)
        db.session.add(new_campaign)
        db.session.commit()
        campaignId = new_campaign.id
        sid = ReliefConfig.query.filter_by(reliefKey = 'current_campaign_id').first()
        try:
            sid.reliefValue = campaignId
            db.session.commit()
        except:
            db.session.rollback()
        else:
            #  Original:   return redirect(url_for('addcyberkillchainphase'))
            return redirect(url_for('displaycampaign'))
    else:
        return render_template("newcampaign.html", title='New Campaign', form=form)

def reset_phase():
    phaseId = ReliefConfig.query.filter_by(reliefKey='current_phase_id').first()
    try:
        phaseId.reliefValue = None
        db.session.commit()
    except:
        print("[e]  Error updating current phase ID record to None.  Rolling back.")
        db.session.rollback()
        return
    else:
        kcupdate_validation = ReliefConfig.query.filter_by(reliefKey = 'current_phase_id').first()
        print("[i]  Updated current phase ID record to None.  Validation: ", kcupdate_validation)
        return

@app.route('/addcyberkillchainphase', methods=["GET", "POST"])
def addcyberkillchainphase():
    campaignId = ReliefConfig.query.filter_by(reliefKey='current_campaign_id').first()
    phaseId = ReliefConfig.query.filter_by(reliefKey='current_phase_id').first()
    form = AddKillChainPhase()
    # Set the form Campaign Name.
    # Use campaignId to pull value from Campaigns table.
    campaignDetails = Campaigns.query.filter_by(id = campaignId.reliefValue).first()
    form.campaignName.data = campaignDetails.campaign_name
    # Use campaignId to pull list of existing phases from the CyberKillChains table
    phases = CyberKillChains.query.filter_by(campaign_id = campaignId.reliefValue).all()
    #print("[d]  KC existing kill chain phases: ", phases)
    # check if current_phase_id set in reliefConfig
    if phaseId.reliefValue is not None:
        print("[d]  Phase ID is not None:", phaseId.reliefValue)
        # Get the phase details
        current_phase_details = CyberKillChains.query.filter_by(id = phaseId.reliefValue).first()
        print("[d]  Details: ", current_phase_details)
    #---
    if 'edit' in request.form:
        # Editing Phase
        print("[d]  Editing KC phase")
        # Create an instance of the add phase form
        form = AddKillChainPhase()
        phaserowid = request.form['editphase_row_id']
        # Update current_phase_id in reliefConfig table
        #  Redundant from above, should probably move here.  sid = ReliefConfig.query.filter_by(reliefKey = 'current_phase_id').first()
        print("[d]  current phase id record value is ", phaseId.reliefValue)
        try:
            phaseId.reliefValue = phaserowid
            db.session.commit()
        except:
            print("[e]  Error updating current phase ID record.  Rolling back.")
            db.session.rollback()
        else:
            # verify
            #kcupdate_validation = ReliefConfig.query.filter_by(current_phase_id = phaseId.reliefValue, phase = form.phase.data).scalar() is not None
            kcupdate_validation = ReliefConfig.query.filter_by(reliefKey = 'current_phase_id').first()
            print("[i]  Updated current phase ID record to ", phaserowid, ".  Validation: ", kcupdate_validation)
            return render_template('addcyberkillchainphase.html', title='Update Phase', form=form)

    elif 'delete' in request.form:
        # Delete the Phase
        form = AddKillChainPhase()
        rowid = request.form['deletephase_row_id']
        print("[d]  Deleting phase id record value is ", rowid)
        current_phase_details = CyberKillChains.query.filter_by(id = rowid).first()
        try:
            db.session.delete(current_phase_details)
            db.session.commit()
        except:
            print("[e]  Error deleting current phase ID record.  Rolling back.")
            db.session.rollback()
        else:
            reset_phase()
            return render_template('addcyberkillchainphase.html', form=form)

    elif 'save' in request.form:
        #----
        #   Original functioning code for add new phase 
        #   prior to adding button handlers above.
        if form.validate_on_submit():
            print("[d]  KC Validating on Submit - Add cyber kill chain phase to the campaign")
            # New phase
            # First, check to see if the phase exists for the campaign
            print("[d]  KC Phase to add ", form.phase.data)
            kc_exist_check = CyberKillChains.query.filter_by(campaign_id = campaignId.reliefValue, phase = form.phase.data).scalar() is not None
            print("[d]  KC does this phase exist for the campaign?  ", kc_exist_check)

            # Convert the multiselectfield list to a csv. 
            # Otherwise there is a data type mismatch with the DB table.
            actionsCsv = ",".join(map(str, form.courses_of_action.data))
            print("[d]  KC actionsCsv: ", actionsCsv)

            # Second, construct the entity and insert (add) to the KC table
            new_kcphase = CyberKillChains(
                campaign_id = campaignId.reliefValue,
                phase = form.phase.data,
                adversary = form.adversary.data,
                capability = form.capability.data,
                infrastructure = form.infrastructure.data,
                victim = form.victim.data,
                business_impact = form.business_impact.data,
                courses_of_action = actionsCsv
            )

            if kc_exist_check == False:
                # Phase does not exist: Insert new record
                print("[d]  KC Insert record", new_kcphase)
                try:
                    db.session.add(new_kcphase)
                    db.session.commit()
                except:
                    flash("Failed to add the Cyber Kill Chain phase to the campaign.  Does the campaign already have this phase?", "danger")
                    db.session.rollback()
                else:
                    flash("Cyber Kill Chain phase added successfully", "success")
            else:
                # Phase exist: Update existing record
                print("[d]  KC Update record")
                try:
                    db.session.commit()
                except:
                    flash("Failed to add the Cyber Kill Chain phase to the campaign.  Does the campaign already have this phase?", "danger")
                    db.session.rollback()
                else:
                    flash("Cyber Kill Chain phase added successfully", "success")
            # Trust but verify.  Did the transaction succeed?
            kc_validation = CyberKillChains.query.filter_by(campaign_id = campaignId.reliefValue, phase = form.phase.data).scalar() is not None
            print("[d]  KC After DML - does this phase exist for the campaign?  ", kc_validation)
            # Refresh the list of phases as an indicator to the user.
            phases = CyberKillChains.query.filter_by(campaign_id = campaignId.reliefValue).all()
            print("[d]  KC existing kill chain phases: ", phases)
    # Reload to add more phases.
    return render_template("addcyberkillchainphase.html", title='Add Kill Chain Phase', phases=phases, form=form)


@app.route('/displaycampaign/', methods=["GET", "POST"])
def displaycampaign():
    # View or update campaign level information or
    # Add or edit kill chain phases
    form = DisplayCampaign()
    rowid = ReliefConfig.query.filter_by(reliefKey = 'current_campaign_id').first().reliefValue
    phaserowid = ReliefConfig.query.filter_by(reliefKey = 'current_phase_id').first().reliefValue
    phases = CyberKillChains.query.filter_by(campaign_id = rowid).all()
    # set form to latest values
    setCampaign = Campaigns.query.filter_by(id = rowid).first()
    form.campaignName.data = setCampaign.campaign_name
    form.reportDate.data = setCampaign.report_date
    form.summaryBusinessImpact.data = setCampaign.summary_of_business_impact

    # Check the button pressed
    if request.method == 'POST':
        if 'savecampaign' in request.form:
            modifyCampaign = Campaigns.query.filter_by(id = rowid).first()
            modifyCampaign.campaign_name = request.form['campaignName']
            #  Currently, reportDate cannot be updated.
            # print("\n\n[d]  Initial state of record: ", setCampaign)
            # tmpReportDate = request.form['reportDate'].split("/")
            # print("[d]  Temp Date List  0:Month:", tmpReportDate[0], "  1:Day:", tmpReportDate[1], "    2:Year:", tmpReportDate[2])
            # reportDateForSql = tmpReportDate[2]+"-"+tmpReportDate[0]+"-"+tmpReportDate[1]+" 00:00:00"   #.000000"
            # print("[d]  Formatted date string: ", reportDateForSql)
            # #modifyCampaign.report_date = reportDateForSql
            # print("[d]  Saving record          : ", modifyCampaign)
            #---<<<---End Report Date troubleshooting.
            modifyCampaign.summary_of_business_impact = request.form['summaryBusinessImpact']
            try:
                db.session.commit()
            except:
                flash("Failed to update campaign data", "danger")
                db.session.rollback()
            else:
                flash("Campaign successfully updated", "success")
                # Refresh form fields with updated values
                setCampaign = Campaigns.query.filter_by(id = rowid).first()
                form.campaignName.data = setCampaign.campaign_name
                form.reportDate.data = setCampaign.report_date
                form.summaryBusinessImpact.data = setCampaign.summary_of_business_impact
                print("[d]  Final record as s: ", setCampaign, "\n\n")

        elif 'deletecampaign' in request.form:
            if rowid is None:
                flash("Failed to delete a non-existent campaign.", "danger")
                return
            else:
                campaignName = Campaigns.query.filter_by(id = rowid).first()
                try:
                    db.session.delete(campaignName)
                    db.session.commit()
                except:
                    flash("Failed to delete the campaign.", "danger")
                    db.session.rollback()
                else:
                    flash("Successfully deleted the campaign.", "success")
            items = Campaigns.query.all()
            form = ListCampaigns()
            return render_template("listcampaigns.html", items=items, form=form)

        elif 'deletephase' in request.form:
            # Delete the phase
            targetPhaseId = CyberKillChains.query.filter_by(id = request.form['deletephase_row_id']).first()
            if targetPhaseId is None:
                flash("Failed to delete a non-existent phase target.", "danger")
            else:
                try:
                    db.session.delete(targetPhaseId)
                    db.session.commit()
                except:
                    flash("Failed to delete the phase.", "danger")
                    db.session.rollback()
                else:
                    flash("Successfully deleted the phase.", "success")
            phases = CyberKillChains.query.filter_by(campaign_id = rowid).all()
                
        elif 'editphase' in request.form:
            # Edit the phase
            print("[d]  Edit phase ", phaserowid)
            phaserowid = request.form['editphase_row_id']
            try:
                db.session.commit()
            except:
                flash("Failed to mark the phase for edit", "danger")
                db.session.rollback()
            else:
                flash("Cyber Kill Chain phase added successfully to the add or modify fields below.", "success")
            # Get phase row details to populate the form.
            oldPhaseDetails = CyberKillChains.query.filter_by(id = request.form['editphase_row_id']).first()
            form.phase.data = oldPhaseDetails.phase
            form.adversary.data = oldPhaseDetails.adversary
            form.capability.data = oldPhaseDetails.capability
            form.infrastructure.data = oldPhaseDetails.infrastructure
            form.victim.data = oldPhaseDetails.victim
            form.business_impact.data = oldPhaseDetails.business_impact
            #  May need extra tweak to update multiselect list.
            form.courses_of_action.data = oldPhaseDetails.courses_of_action

        elif 'addphase' in request.form:
            actionsCsv = ",".join(map(str, form.courses_of_action.data))
            kc_exist_check = CyberKillChains.query.filter_by(campaign_id = rowid, phase = form.phase.data).scalar() is not None
            if kc_exist_check == False:
                # Phase does not exist: Insert new record
                new_kcphase = CyberKillChains(
                    campaign_id = rowid,
                    phase = form.phase.data,
                    adversary = form.adversary.data,
                    capability = form.capability.data,
                    infrastructure = form.infrastructure.data,
                    victim = form.victim.data,
                    business_impact = form.business_impact.data,
                    courses_of_action = actionsCsv
                )
                try:
                    db.session.add(new_kcphase)
                    db.session.commit()
                except:
                    flash("Failed to add the Cyber Kill Chain phase to the campaign.  Does the campaign already have this phase?", "danger")
                    db.session.rollback()
                else:
                    flash("Cyber Kill Chain phase added successfully", "success")
            else:
                # Phase exist: Update existing record
                # Get the current phase row data to update
                mod_kcphase = CyberKillChains.query.filter_by(campaign_id = rowid, phase = form.phase.data).first()
                # Update record values - original idea for statement
                mod_kcphase.campaign_id = rowid
                mod_kcphase.phase = form.phase.data
                mod_kcphase.adversary = form.adversary.data
                mod_kcphase.capability = form.capability.data
                mod_kcphase.infrastructure = form.infrastructure.data
                mod_kcphase.victim = form.victim.data
                mod_kcphase.business_impact = form.business_impact.data
                mod_kcphase.courses_of_action = actionsCsv
                try:
                    db.session.commit()
                except:
                    flash("Failed to add the Cyber Kill Chain phase to the campaign.  Does the campaign already have this phase?", "danger")
                    db.session.rollback()
                else:
                    flash("Cyber Kill Chain phase added successfully", "success")
                    verify = CyberKillChains.query.filter_by(campaign_id = rowid).first()
            # Refresh the phases table for the form
            phases = CyberKillChains.query.filter_by(campaign_id = rowid).all()
    elif request.method == 'GET':
        print("[d]  Setting form to current campaign row values.")
    phases = CyberKillChains.query.filter_by(campaign_id = rowid).all()
    return render_template("displaycampaign.html", title='Enrich Campaign', phases=phases, form=form)

@app.route('/listcampaigns/', methods=["GET", "POST"])
def listcampaigns():
    form = ListCampaigns()
    items = Campaigns.query.all()
    
    if form.validate_on_submit():
        rowid = request.form['enrich_row_id']
        campaignName = Campaigns.query.filter_by(id = rowid).first()
    
    # Check the button pressed
    if request.method == 'POST':
        # Make Presentation
        if 'export' in request.form:
            rowid = form.export_row_id.data
            try:
                Make_Presentation(rowid)
            except:
                flash("Failed to create presentation in the export directory", "danger")
            else:
                flash("Presentation successfully created in the export directory", "success")
            try:
                Make_Report(rowid)
            except:
                flash("Failed to create report in the export directory", "danger")
            else:
                flash("Report successfully created in the export directory", "success")

        # Enrich Campaign
        elif 'enrich' in request.form:
            rowid = request.form['enrich_row_id']
            #campaignName = Campaigns.query.filter_by(id = rowid).first()
            if rowid == '':
                flash("Failed to enrich the campaign.", "danger")
                return
            else:
                #phases = CyberKillChains.query.filter_by(campaign_id = rowid).all()
                # Update "current campaign ID stored in DB setting".
                sid = ReliefConfig.query.filter_by(reliefKey = 'current_campaign_id').first()
                try:
                    sid.reliefValue = rowid
                    db.session.commit()
                except:
                    flash("Failed to enrich the campaign.", "danger")
                    db.session.rollback()
            return redirect(url_for('displaycampaign'))
            
        # Delete campaign
        elif 'delete' in request.form:
            rowid = request.form['delete_row_id']
            if rowid == '':
                flash("Failed to delete the campaign.", "danger")
                return
            else:
                campaignName = Campaigns.query.filter_by(id = rowid).first()
                try:
                    db.session.delete(campaignName)
                    db.session.commit()
                except:
                    flash("Failed to delete the campaign.", "danger")
                    db.session.rollback()
                else:
                    flash("Successfully deleted the campaign.", "success")
                    items = Campaigns.query.all()
        
    elif request.method == 'GET':
        return render_template("listcampaigns.html", items=items, form=form)
    return render_template("listcampaigns.html", items=items, form=form)

def Make_Report(campaignId):
    print("\n\n\n[d]  Creating report for campaignId ", campaignId)
    killChainNames = {"Recon": 'Reconnaissance', "Weapon": "Weaponization", "Delivery": "Delivery", "Exploit": "Exploitation", "Install": "Installation", "C2": "Command and Control", "Actions": "Actions on Objectives"}
    chains = CyberKillChains.query.filter_by(campaign_id = campaignId).all()
    # current working directory
    dir_path = os.getcwd()
    # configuration settings
    tlpImage = {"RED": dir_path+'/relief/static/images/tlp-red.png', 
        "AMBER": dir_path+'/relief/static/images/tlp-amber.png', 
        "GREEN": dir_path+'/relief/static/images/tlp-green.png', 
        "WHITE": dir_path+'/relief/static/images/tlp-white.png'}
    filenamePrefix = ReliefConfig.query.filter_by(reliefKey = 'filenamePrefix').first().reliefValue
    dataMarkingFooter = ReliefConfig.query.filter_by(reliefKey = 'dataMarkingFooter').first().reliefValue
    presentationAuthor = ReliefConfig.query.filter_by(reliefKey = 'presentationAuthor').first().reliefValue
    presentationCategory = ReliefConfig.query.filter_by(reliefKey = 'presentationCategory').first().reliefValue
    presentationKeywords = ReliefConfig.query.filter_by(reliefKey = 'presentationKeywords').first().reliefValue
    presentationSubject = ReliefConfig.query.filter_by(reliefKey = 'presentationSubject').first().reliefValue
    presentationTitle = ReliefConfig.query.filter_by(reliefKey = 'presentationTitle').first().reliefValue
    trafficLightDefault = ReliefConfig.query.filter_by(reliefKey = 'trafficLightDefault').first().reliefValue
    # Get campaign data
    campaignName = Campaigns.query.filter_by(id = campaignId).first().campaign_name
    origDate = Campaigns.query.filter_by(id = campaignId).first().report_date
    origDate = str(origDate)
    d = datetime.datetime.strptime(origDate, '%Y-%m-%d %H:%M:%S')
    reportDate = d.strftime('%m/%d/%y')
    summaryBusinessImpact = Campaigns.query.filter_by(id = campaignId).first().summary_of_business_impact
    # Get Cyber Kill Chain phases for the campaign
    #phases = CyberKillChains.query.filter_by(campaignId = campaignId).all()
    # Set the presenation output filename
    presentationExportFileName = dir_path + "/relief/export/" + filenamePrefix + "-" + campaignName.replace(" ", "_") + "-" + reportDate.replace("/", "-") + ".docx"
    # Instantiate a new document
    document = Document()
    # Header
    section = document.sections[0]
    header = section.header
    header
    header.is_linked_to_previous
    paragraph = header.paragraphs[0]
    paragraph.text = campaignName + "\t" + reportDate + "\t"
    paragraph.style = document.styles["Header"]
    r = paragraph.add_run()
    r.add_picture(tlpImage[trafficLightDefault])
    # Footer
    footerSection = document.sections[0]
    footer = footerSection.footer
    footer
    footer.is_linked_to_previous
    footer_p0 = footer.paragraphs[0]
    footer_p0.text = "\t" + dataMarkingFooter + "\t"
    footer_p0.style = document.styles["Footer"]
    rfooter = footer_p0.add_run()
    rfooter.add_picture(tlpImage[trafficLightDefault])
    # Title page
    document.add_heading(campaignName, 0)
    p0 = document.add_paragraph("by " + presentationAuthor)
    p1 = document.add_paragraph("Report Date: " + reportDate + "\t")
    document.add_page_break()
    # Executive Summary
    document.add_heading('Executive Summary', level=1)
    p0 = document.add_paragraph(summaryBusinessImpact)
    p1 = document.add_paragraph("Briefly describe the root cause of the event.")
    p2 = document.add_paragraph("How are you able to detect and deny future events based on that knowledge?")
    p3 = document.add_paragraph("What are several positive outcomes from those defensive actions?  For example, how many other times has this occurred since and been stopped?")
    # idea to have each phase's summary with a bubble for its phase.  Picture API isn't mature(?)
    document.add_heading('Cyber Kill Chain Phase Overview', level=1)
    for chain in chains:
        # Cyber Kill Chain phase summary bubbles
        document.add_heading(chain.phase, level=2)
        p4 = document.add_paragraph(chain.business_impact)
        #document.add_picture(dir_path + '/relief/static/images/phase-disabled.png', width=Inches(1.0))
    document.add_page_break()
    # Cyber Kill Chain pages
    for chain in chains:
        # Overview of Cyber Kill Chain phases
        document.add_heading(chain.phase, level=1)
        document.add_heading('Business Impact', level=2)
        p0 = document.add_paragraph(chain.business_impact)
        document.add_heading('Courses of Action', level=2)
        p0 = document.add_paragraph(chain.courses_of_action)
        document.add_heading('Lessons Learned', level=2)
        p0 = document.add_paragraph('')
        document.add_heading('Adversary', level=2)
        p0 = document.add_paragraph(chain.adversary)
        document.add_heading('Capability', level=2)
        p0 = document.add_paragraph(chain.capability)
        document.add_heading('Infrastructure', level=2)
        p0 = document.add_paragraph(chain.infrastructure)
        document.add_heading('Victim', level=2)
        p0 = document.add_paragraph(chain.victim)
        document.add_page_break()
    # Write the document to a file
    document.save(presentationExportFileName)
    return

def Make_Presentation(campaignId):
    # Instantiate the presentation deck
    prs = Presentation()
    # Abbreviated cyber kill chain names for menu links on executive summary slide
    killChainNames = {"Recon": 'Reconnaissance', "Weapon": "Weaponization", "Delivery": "Delivery", "Exploit": "Exploitation", "Install": "Installation", "C2": "Command and Control", "Actions": "Actions on Objectives"}
    #killChainNames = {'Reconnaissance': "Recon", "Weaponization":"Weapon", "Delivery": "Delivery", "Exploitation":"Exploit", "Installation":"Install", "Command and Control":"C2", "Actions on Objectives":"Actions"}
    # cyber kill chain phase rows for this campaign
    chains = CyberKillChains.query.filter_by(campaign_id = campaignId).all()
    chains_for_matching = ""
    for chain in chains:
        chains_for_matching += " " + ''.join(chain.phase)
    # current working directory
    dir_path = os.getcwd()
    # traffic light protocol
    tlpImage = {"RED": dir_path+'/relief/static/images/tlp-red.png', 
        "AMBER": dir_path+'/relief/static/images/tlp-amber.png', 
        "GREEN": dir_path+'/relief/static/images/tlp-green.png', 
        "WHITE": dir_path+'/relief/static/images/tlp-white.png'}
    # configuration settings
    filenamePrefix = ReliefConfig.query.filter_by(reliefKey = 'filenamePrefix').first().reliefValue
    dataMarkingFooter = ReliefConfig.query.filter_by(reliefKey = 'dataMarkingFooter').first().reliefValue
    presentationAuthor = ReliefConfig.query.filter_by(reliefKey = 'presentationAuthor').first().reliefValue
    presentationCategory = ReliefConfig.query.filter_by(reliefKey = 'presentationCategory').first().reliefValue
    presentationKeywords = ReliefConfig.query.filter_by(reliefKey = 'presentationKeywords').first().reliefValue
    presentationSubject = ReliefConfig.query.filter_by(reliefKey = 'presentationSubject').first().reliefValue
    presentationTitle = ReliefConfig.query.filter_by(reliefKey = 'presentationTitle').first().reliefValue
    trafficLightDefault = ReliefConfig.query.filter_by(reliefKey = 'trafficLightDefault').first().reliefValue
    # Get campaign data
    campaignName = Campaigns.query.filter_by(id = campaignId).first().campaign_name
    origDate = Campaigns.query.filter_by(id = campaignId).first().report_date
    origDate = str(origDate)
    d = datetime.datetime.strptime(origDate, '%Y-%m-%d %H:%M:%S')
    reportDate = d.strftime('%m/%d/%y')
    summaryBusinessImpact = Campaigns.query.filter_by(id = campaignId).first().summary_of_business_impact
    # Set the presenation output filename
    presentationExportFileName = dir_path + "/relief/export/" + filenamePrefix + "-" + campaignName.replace(" ", "_") + "-" + reportDate.replace("/", "-") + ".pptx"

    #----
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    # Title slide content
    title.text = campaignName
    subtitle.text = reportDate
    # Every slide gets the TLP. Orient the traffic light right top & bottom
    # Top Right
    left = Inches(8.5)
    top = Inches(0.25)
    tlpTop = slide.shapes.add_picture(tlpImage[trafficLightDefault], left, top)
    # Bottom Right
    top = Inches(7)
    tlpBottom = slide.shapes.add_picture(tlpImage[trafficLightDefault], left, top)
    # Marking and handling classification
    left = Inches(3.5)
    top = Inches(6.75)
    height = Inches(0.25)
    width = Inches(3.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = dataMarkingFooter
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(10)
    
    #----
    # Executive Overview
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    # Every slide gets the TLP. Orient the traffic light right top & bottom
    # Top Right
    left = Inches(8.5)
    top = Inches(0.25)
    tlpTop = slide.shapes.add_picture(tlpImage[trafficLightDefault], left, top)
    # Bottom Right
    top = Inches(7)
    tlpBottom = slide.shapes.add_picture(tlpImage[trafficLightDefault], left, top)
    bullet_slide_layout = prs.slide_layouts[1]
    # Marking and handling classification
    left = Inches(3.5)
    top = Inches(6.75)
    height = Inches(0.25)
    width = Inches(3.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = dataMarkingFooter
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(10)
    
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Executive Summary'

    tf = body_shape.text_frame
    p0 = tf.add_paragraph()
    p0.text = summaryBusinessImpact
    p1 = tf.add_paragraph()
    p1.text = 'Campaign Name: ' + campaignName


    # Kill chain hotlinks
    left = Inches(0.7)  # 0.93" centers this overall set of shapes
    top = Inches(6.0)
    width = Inches(1.25)
    height = Inches(.5)
    for n in killChainNames:
        chainPhaseShape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        chainPhaseShape.text = '%s' % n #killChainNames[n]
        left = left + width - Inches(0.0)
        line = chainPhaseShape.line
        line.width = Pt(3.0)
        line.color.rgb = RGBColor(147, 74, 187)
        # if completed KC phase...color orange...else gray
        if killChainNames[n] in chains_for_matching:
            fill = chainPhaseShape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 123, 89)
        else:
            fill = chainPhaseShape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(221, 221, 221)


    #----
    # Cyber Kill Chain & Diamond Model detail slides
    for chain in chains:
        #print("[d]  chain ::: ", chain)
        #----
        # Business Impact for Phase
        # reset values
        left = Inches(0)
        top = Inches(0)
        width = Inches(0)
        height = Inches(0)

        phase_summary = prs.slide_layouts[5]
        phase_summary_slide = prs.slides.add_slide(phase_summary)
        shapes = phase_summary_slide.shapes

        # Every slide gets the TLP. Orient the traffic light right top & bottom
        #  A way to do this with Python-PPTX without images: (future) 
        #       font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)
        # Orient the traffic light right top & bottom
        # Top Right
        left = Inches(8.5)
        top = Inches(0.25)
        tlpTop = phase_summary_slide.shapes.add_picture(tlpImage[trafficLightDefault], left, top)
        # Bottom Right
        top = Inches(7)
        tlpBottom = phase_summary_slide.shapes.add_picture(tlpImage[trafficLightDefault], left, top)

        # Marking and handling classification
        left = Inches(3.5)
        top = Inches(6.75)
        height = Inches(0.25)
        width = Inches(3.0)
        txBox = phase_summary_slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = dataMarkingFooter
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(10)
        
        # Title
        left = Inches(0.0)
        top = Inches(0.0)
        shapes.title.text = chain.phase


        # Business Impact content
        left = Inches(1.0)
        top = Inches(2.0)
        width = Inches(8.0)
        txBox = phase_summary_slide.shapes.add_textbox(left, top, width, height)
        line = txBox.line
        line.color.rgb = RGBColor(0,0,0)
        line.width = Pt(1.0)
        tf = txBox.text_frame
        tf.word_wrap = True
        p1 = tf.add_paragraph()
        p1.text = "Business Impact"
        p1.alignment = PP_ALIGN.CENTER
        p1.font.bold = True
        p1.font.size = Pt(18)
        p1.font.color.rgb = RGBColor(183, 43, 43)
        p0 = tf.add_paragraph()
        p0.text = chain.business_impact
        p0.alignment = PP_ALIGN.CENTER
        p0.font.color.rgb = RGBColor(183, 43, 43)
        
        # Courses of Action content
        left = Inches(1.0)
        top = Inches(5.0)
        width = Inches(8.0)
        txBox = phase_summary_slide.shapes.add_textbox(left, top, width, height)
        line = txBox.line
        line.color.rgb = RGBColor(0,0,0)
        line.width = Pt(1.0)
        tf = txBox.text_frame
        tf.word_wrap = True
        p1 = tf.add_paragraph()
        p1.text = "Courses of Action"
        p1.alignment = PP_ALIGN.CENTER
        p1.font.bold = True
        p1.font.size = Pt(18)
        p1.font.color.rgb = RGBColor(43, 127, 183)
        p0 = tf.add_paragraph()
        p0.text = chain.courses_of_action
        p0.alignment = PP_ALIGN.CENTER
        p0.font.color.rgb = RGBColor(43, 127, 183)


        #---
        #  New slide for diamond model details
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        # Every slide gets the TLP. Orient the traffic light right top & bottom
        #  A way to do this with Python-PPTX without images: (future) 
        #       font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)
        # Orient the traffic light right top & bottom
        # Top Right
        left = Inches(8.5)
        top = Inches(0.25)
        tlpTop = slide.shapes.add_picture(tlpImage[trafficLightDefault], left, top)
        # Bottom Right
        top = Inches(7)
        tlpBottom = slide.shapes.add_picture(tlpImage[trafficLightDefault], left, top)

        # Marking and handling classification
        left = Inches(3.5)
        top = Inches(6.75)
        height = Inches(0.25)
        width = Inches(3.0)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = dataMarkingFooter
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(10)
        
        # Title
        left = Inches(0.0)
        top = Inches(0.0)
        shapes.title.text = chain.phase

        # Diamond image placement
        left = Inches(6.25)  # 0.93" centers this overall set of shapes
        top = Inches(3.0)
        width = Inches(1.25)
        height = Inches(1.5)
        diamond = shapes.add_shape(MSO_SHAPE.DIAMOND, left, top, width, height)
        fill = diamond.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 123, 89)
        line = diamond.line
        line.width = Pt(6.0)
        line.color.rgb = RGBColor(147, 74, 187)

        # Business Impact content
        left = Inches(0.2)
        top = Inches(1.25)
        width = Inches(3.0)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        line = txBox.line
        line.color.rgb = RGBColor(0,0,0)
        line.width = Pt(1.0)
        tf = txBox.text_frame
        tf.word_wrap = True
        p1 = tf.add_paragraph()
        p1.text = "Business Impact"
        p1.alignment = PP_ALIGN.CENTER
        p1.font.bold = True
        p1.font.size = Pt(18)
        p1.font.color.rgb = RGBColor(183, 43, 43)
        p0 = tf.add_paragraph()
        p0.text = chain.business_impact
        p0.alignment = PP_ALIGN.CENTER
        p0.font.color.rgb = RGBColor(183, 43, 43)
        
        # Courses of Action content
        left = Inches(0.2)
        top = Inches(5.25)
        width = Inches(3.0)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        line = txBox.line
        line.color.rgb = RGBColor(0,0,0)
        line.width = Pt(1.0)
        tf = txBox.text_frame
        tf.word_wrap = True
        p1 = tf.add_paragraph()
        p1.text = "Courses of Action"
        p1.alignment = PP_ALIGN.CENTER
        p1.font.bold = True
        p1.font.size = Pt(18)
        p1.font.color.rgb = RGBColor(43, 127, 183)
        p0 = tf.add_paragraph()
        p0.text = chain.courses_of_action
        p0.alignment = PP_ALIGN.CENTER
        p0.font.color.rgb = RGBColor(43, 127, 183)

        # Adversary content
        left = Inches(5.5)
        top = Inches(1.75)
        width = Inches(3.0)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p0 = tf.add_paragraph()
        p0.text = chain.adversary
        p0.alignment = PP_ALIGN.CENTER
        p0.font.bold = False
        p1 = tf.add_paragraph()
        p1.text = "Adversary"
        p1.alignment = PP_ALIGN.CENTER
        p1.font.bold = True

        # Victim content
        left = Inches(5.5)
        top = Inches(4.75)
        width = Inches(3.0)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p0 = tf.add_paragraph()
        p0.text = chain.victim
        p0.alignment = PP_ALIGN.CENTER
        p0.font.bold = False
        p1 = tf.add_paragraph()
        p1.text = "Victim"
        p1.alignment = PP_ALIGN.CENTER
        p1.font.bold = True

        # Capability content
        left = Inches(7.85)
        top = Inches(3.0)
        width = Inches(2.0)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p0 = tf.add_paragraph()
        p0.text = chain.capability
        p0.alignment = PP_ALIGN.CENTER
        p0.font.bold = False
        p1 = tf.add_paragraph()
        p1.text = "Capability"
        p1.alignment = PP_ALIGN.CENTER
        p1.font.bold = True

        # Infrastructure content
        left = Inches(3.0)
        top = Inches(3.0)
        width = Inches(3.0)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p0 = tf.add_paragraph()
        p0.text = chain.infrastructure
        p0.alignment = PP_ALIGN.CENTER
        p0.font.bold = False
        p1 = tf.add_paragraph()
        p1.text = "Infrastructure"
        p1.alignment = PP_ALIGN.CENTER
        p1.font.bold = True
        #----End Cyber Kill Chain phase loop
    #----
    # Write the slides to a new presentation file
    prs.save(presentationExportFileName)
    return


@app.route('/updatesettingsform/', methods=["GET", "POST"])
def upset():
    form = UpdateSettingsForm()
    #  There is only one Submit button may need to extend the widget.
    #  Leaving for now.
    #if form.process_button.data == UpdateSettingsForm.RESET:
    #    print("[d]   Caught me a reset request")
    if form.validate_on_submit():
        s = ReliefConfig.query.filter_by(reliefKey = 'filenamePrefix').first()
        s.reliefValue = form.newFileNamePrefix.data
        db.session.commit()

        s = ReliefConfig.query.filter_by(reliefKey = 'trafficLightDefault').first()
        s.reliefValue = form.newTrafficLightDefault.data
        db.session.commit()
        
        s = ReliefConfig.query.filter_by(reliefKey = 'dataMarkingFooter').first()
        s.reliefValue = form.newDataMarkingFooter.data
        db.session.commit()
                
        s = ReliefConfig.query.filter_by(reliefKey = 'presentationAuthor').first()
        s.reliefValue = form.newPresentationAuthor.data
        db.session.commit()
        
        s = ReliefConfig.query.filter_by(reliefKey = 'presentationCategory').first()
        s.reliefValue = form.newPresentationCategory.data
        db.session.commit()
        
        s = ReliefConfig.query.filter_by(reliefKey = 'presentationTitle').first()
        s.reliefValue = form.newPresentationTitle.data
        db.session.commit()
        
        s = ReliefConfig.query.filter_by(reliefKey = 'presentationSubject').first()
        s.reliefValue = form.newPresentationSubject.data
        db.session.commit()
        
        s = ReliefConfig.query.filter_by(reliefKey = 'presentationKeywords').first()
        s.reliefValue = form.newPresentationKeywords.data
        db.session.commit()

        flash("Settings updated successfully", "success")
        print("Settings updated successfully", "success")
    elif request.method == 'GET':
        print("Got the Get")
        s = ReliefConfig.query.filter_by(reliefKey = 'filenamePrefix').first()
        form.newFileNamePrefix.data = s.reliefValue
        
        s = ReliefConfig.query.filter_by(reliefKey = 'trafficLightDefault').first()
        form.newTrafficLightDefault.data = s.reliefValue
        
        s = ReliefConfig.query.filter_by(reliefKey = 'dataMarkingFooter').first()
        form.newDataMarkingFooter.data = s.reliefValue
                
        s = ReliefConfig.query.filter_by(reliefKey = 'presentationAuthor').first()
        form.newPresentationAuthor.data = s.reliefValue
        
        s = ReliefConfig.query.filter_by(reliefKey = 'presentationCategory').first()
        form.newPresentationCategory.data = s.reliefValue
        
        s = ReliefConfig.query.filter_by(reliefKey = 'presentationTitle').first()
        form.newPresentationTitle.data = s.reliefValue
        
        s = ReliefConfig.query.filter_by(reliefKey = 'presentationSubject').first()
        form.newPresentationSubject.data = s.reliefValue
        
        s = ReliefConfig.query.filter_by(reliefKey = 'presentationKeywords').first()
        form.newPresentationKeywords.data = s.reliefValue
        
    return render_template("updatesettingsform.html", title='Update Settings Form', form=form)